import collections
import collections.abc
import os
from pptx import Presentation
from pptx.util import Inches, Pt

def add_slide(prs, image_path, title_text, answer_text, explanation_text):
    # --- SLIDE 1: SOAL SAJA (Efek Animasi) ---
    slide_layout = prs.slide_layouts[5] # blank slide or title slide
    slide1 = prs.slides.add_slide(slide_layout)
    
    # Set Title
    title1 = slide1.shapes.title
    title1.text = title_text
    
    # Add Image at the top, centered
    left = Inches(1)
    top = Inches(1.2)
    width = Inches(8)
    
    try:
        pic1 = slide1.shapes.add_picture(image_path, left, top, width=width)
    except Exception as e:
        print(f"Could not add image {image_path}: {e}")
        
    # --- SLIDE 2: SOAL + JAWABAN ---
    slide2 = prs.slides.add_slide(slide_layout)
    
    # Set Title
    title2 = slide2.shapes.title
    title2.text = title_text
    
    try:
        pic2 = slide2.shapes.add_picture(image_path, left, top, width=width)
    except:
        pass
        
    # Add Textbox for Answer & Explanation at the bottom
    left_tx = Inches(1)
    top_tx = Inches(4.5)  # Positioned below the image
    width_tx = Inches(8)
    height_tx = Inches(2.8)
    
    txBox = slide2.shapes.add_textbox(left_tx, top_tx, width_tx, height_tx)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Jawaban:"
    p.font.bold = True
    p.font.size = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = answer_text
    p2.font.size = Pt(18)
    
    p3 = tf.add_paragraph()
    p3.text = "\nPenjelasan:"
    p3.font.bold = True
    p3.font.size = Pt(20)
    
    p4 = tf.add_paragraph()
    p4.text = explanation_text
    p4.font.size = Pt(16)


def main():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Pembahasan Soal Sertifikasi"
    subtitle.text = "Bagian 1 (Soal 1 - 5)"

    data = [
        {
            "image": "Sample Soft Dev/1.jpg",
            "title": "Soal 1: Struktur Data Stack",
            "answer": "C. 6",
            "explanation": "Awal: [2, 4, 6, 8] (Top = 8)\nPop -> [2, 4, 6]\nPush 3 -> [2, 4, 6, 3]\nPop -> [2, 4, 6]\nPush 4 -> [2, 4, 6, 4]\nPush 6 -> [2, 4, 6, 4, 6]\nPush 7 -> [2, 4, 6, 4, 6, 7]\nPop -> [2, 4, 6, 4, 6]\nPop -> [2, 4, 6, 4]\nPop -> [2, 4, 6] (Top = 6)"
        },
        {
            "image": "Sample Soft Dev/2.jpg",
            "title": "Soal 2: C# Access Modifier",
            "answer": "C. protected",
            "explanation": "'protected' membatasi akses hanya di dalam class itu sendiri dan class yang mewarisinya (derived class). 'private' hanya untuk class itu sendiri, 'internal' untuk assembly yang sama, dan 'public' bisa diakses dari mana saja."
        },
        {
            "image": "Sample Soft Dev/3-1.jpg",
            "title": "Soal 3: MVC Architecture",
            "answer": "View -> Controller -> Model -> Controller -> View",
            "explanation": "Alur standar MVC:\n1. View: User clicks button\n2. Controller: Converts event into request\n3. Model: Validates request, stores data\n4. Controller: Interprets and passes data\n5. View: Displays data"
        },
        {
            "image": "Sample Soft Dev/4.jpg",
            "title": "Soal 4: NoSQL Document Database",
            "answer": "C. As XML (atau JSON)",
            "explanation": "Document database (seperti MongoDB, CouchDB) menyimpan data dalam format dokumen terstruktur seperti JSON, BSON, atau XML. Opsi lain merujuk ke tipe database yang berbeda (Nodes = Graph, Rows/Columns = Relational)."
        },
        {
            "image": "Sample Soft Dev/5.jpg",
            "title": "Soal 5: Web Service Serialization",
            "answer": "B. Serialization",
            "explanation": "Serialisasi (Serialization) adalah proses menerjemahkan struktur data atau state dari sebuah objek menjadi format yang dapat disimpan atau ditransmisikan (seperti XML atau JSON). Proses sebaliknya disebut Deserialization."
        }
    ]
    
    base_path = r"c:\Users\PREDATOR\Documents\KULIAH\latihan sertifikasi internasional"
    for item in data:
        img_path = os.path.join(base_path, item["image"])
        add_slide(prs, img_path, item["title"], item["answer"], item["explanation"])
        
    output_file = os.path.join(base_path, "Pembahasan_Soal_V2.pptx")
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

if __name__ == "__main__":
    main()
