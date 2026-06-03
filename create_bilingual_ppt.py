import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def get_images_for_question(q_no, img_dir):
    images = []
    if not os.path.exists(img_dir):
        return images
    for f in sorted(os.listdir(img_dir)):
        if not f.endswith(".jpg"): continue
        name = f[:-4]
        if name == str(q_no) or name.startswith(f"{q_no}-"):
            images.append(os.path.join(img_dir, f))
    return images

def add_slide(prs, item, img_dir, show_answer=False):
    slide_layout = prs.slide_layouts[5] # blank slide with title
    slide = prs.slides.add_slide(slide_layout)
    
    q_no = item["no"]
    title = slide.shapes.title
    title.text = f"Question {q_no} / Soal {q_no}"
    
    images = get_images_for_question(q_no, img_dir)
    
    # Layout dimensions (Widescreen 13.33 x 7.5)
    margin = Inches(0.5)
    
    # If we have images, put them on the left half. If not, text takes full width.
    has_img = len(images) > 0
    text_left = Inches(6.5) if has_img else margin
    text_width = Inches(6.5) if has_img else Inches(12.33)
    
    if has_img:
        # Place images on the left side
        img_top = Inches(1.5)
        for idx, img_path in enumerate(images):
            try:
                # We stack images vertically if multiple
                # Restrict width to max 5.5 inches
                slide.shapes.add_picture(img_path, margin, img_top + Inches(idx*2.5), width=Inches(5.5))
            except Exception as e:
                print(f"Error loading image {img_path}: {e}")

    # English Text Box
    txBox_en = slide.shapes.add_textbox(text_left, Inches(1.5), text_width, Inches(2))
    tf_en = txBox_en.text_frame
    tf_en.word_wrap = True
    
    p = tf_en.add_paragraph()
    p.text = "[English]\n" + item.get("soal_en", "") + "\n\n" + item.get("pilihan_en", "")
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 0, 128) # Navy Blue
    
    # Indonesian Text Box
    txBox_id = slide.shapes.add_textbox(text_left, Inches(4), text_width, Inches(1.5))
    tf_id = txBox_id.text_frame
    tf_id.word_wrap = True
    
    p2 = tf_id.add_paragraph()
    p2.text = "[Indonesian]\n" + item.get("soal_id", "")
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(128, 0, 0) # Dark Red

    if show_answer:
        # Add Answer Box at the bottom
        ans_top = Inches(5.8)
        ans_width = Inches(12.33)
        txBox_ans = slide.shapes.add_textbox(margin, ans_top, ans_width, Inches(1.5))
        
        # Fill background color
        fill = txBox_ans.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(230, 240, 255) # Light Blue Background
        
        tf_ans = txBox_ans.text_frame
        tf_ans.word_wrap = True
        
        p3 = tf_ans.add_paragraph()
        p3.text = f"Answer: {item.get('jawaban', '')}"
        p3.font.bold = True
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(0, 100, 0) # Dark Green
        
        p4 = tf_ans.add_paragraph()
        p4.text = f"Penjelasan: {item.get('penjelasan', '')}"
        p4.font.size = Pt(14)
        p4.font.color.rgb = RGBColor(0, 0, 0)

def main():
    base_dir = r"c:\Users\PREDATOR\Documents\KULIAH\latihan sertifikasi internasional"
    data_file = os.path.join(base_dir, "data.json")
    img_dir = os.path.join(base_dir, "Sample Soft Dev")
    output_file = os.path.join(base_dir, "Software_Development_Exam_QA.pptx")
    
    with open(data_file, 'r', encoding='utf-8') as f:
        data = json.load(f)
        
    prs = Presentation()
    # Set to Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Software Development Exam Q&A"
    subtitle.text = "Bilingual Version (English & Indonesian)\nFull 40 Questions"
    
    for item in data:
        # Slide 1: Question Only
        add_slide(prs, item, img_dir, show_answer=False)
        # Slide 2: Question + Answer
        add_slide(prs, item, img_dir, show_answer=True)
        
    prs.save(output_file)
    print(f"Successfully generated {output_file} with {len(data)} questions.")

if __name__ == "__main__":
    main()
